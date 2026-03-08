"""
SEFS File Reader - Extract text content from various file formats.
Supports: plain text, PDF, DOCX, PPTX, CSV/XLSX, images (OCR), and all programming files.
"""

import csv
import io
import logging
import shutil
import time
from pathlib import Path
from typing import Optional

import chardet

from backend.config import MAX_FILE_SIZE_BYTES

logger = logging.getLogger("sefs.file_reader")

# ─── Binary format extensions (need special readers) ─────────────────────────
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif", ".webp"}

# All extensions we can handle (binary + text from config)
BINARY_EXTENSIONS = PDF_EXTENSIONS | DOCX_EXTENSIONS | PPTX_EXTENSIONS | EXCEL_EXTENSIONS | IMAGE_EXTENSIONS

# Check if tesseract is available for OCR
_tesseract_available: Optional[bool] = None


def _check_tesseract() -> bool:
    """Check if tesseract binary is installed."""
    global _tesseract_available
    if _tesseract_available is None:
        _tesseract_available = shutil.which("tesseract") is not None
        if _tesseract_available:
            logger.info("Tesseract OCR is available")
        else:
            logger.warning("Tesseract OCR not found. Image text extraction disabled. Install with: brew install tesseract")
    return _tesseract_available


def read_text_file(path: Path) -> Optional[str]:
    """Read a plain text file with encoding detection."""
    max_retries = 3
    for attempt in range(max_retries):
        try:
            size = path.stat().st_size
            if size == 0:
                return ""
            if size > MAX_FILE_SIZE_BYTES:
                logger.info(f"Skipping {path} - too large ({size} bytes)")
                return None

            raw = path.read_bytes()
            detected = chardet.detect(raw)
            encoding = detected.get("encoding", "utf-8") or "utf-8"

            try:
                return raw.decode(encoding, errors="replace")
            except (UnicodeDecodeError, LookupError):
                return raw.decode("utf-8", errors="replace")
        except PermissionError:
            if attempt < max_retries - 1:
                time.sleep(0.5)
                continue
            logger.error(f"Permission denied reading {path} after {max_retries} attempts")
            return None
        except Exception as e:
            logger.error(f"Error reading text file {path}: {e}")
            return None


def read_pdf(path: Path) -> Optional[str]:
    """Extract text from a PDF file using PyMuPDF (fitz)."""
    try:
        import fitz  # pymupdf

        doc = fitz.open(str(path))
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()

        text = "\n".join(text_parts).strip()
        if not text:
            logger.info(f"PDF has no extractable text (may be scanned): {path}")
            return _try_ocr_pdf(path)
        return text
    except ImportError:
        logger.warning("pymupdf not installed. Cannot read PDF files.")
        return None
    except Exception as e:
        logger.error(f"Error reading PDF {path}: {e}")
        return None


def _try_ocr_pdf(path: Path) -> Optional[str]:
    """Try OCR on a scanned PDF."""
    if not _check_tesseract():
        return f"[PDF file: {path.name} - no extractable text, OCR unavailable]"
    try:
        import fitz
        from PIL import Image
        import pytesseract

        doc = fitz.open(str(path))
        text_parts = []
        for page_num in range(min(len(doc), 10)):  # Limit to first 10 pages
            page = doc[page_num]
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            page_text = pytesseract.image_to_string(img)
            if page_text.strip():
                text_parts.append(page_text)
        doc.close()
        return "\n".join(text_parts).strip() or None
    except Exception as e:
        logger.error(f"OCR failed for PDF {path}: {e}")
        return None


def read_docx(path: Path) -> Optional[str]:
    """Extract text from a DOCX file."""
    try:
        from docx import Document

        doc = Document(str(path))
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        return "\n".join(text_parts).strip() or None
    except ImportError:
        logger.warning("python-docx not installed. Cannot read DOCX files.")
        return None
    except Exception as e:
        logger.error(f"Error reading DOCX {path}: {e}")
        return None


def read_pptx(path: Path) -> Optional[str]:
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_texts.append(row_text)
            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

        return "\n\n".join(text_parts).strip() or None
    except ImportError:
        logger.warning("python-pptx not installed. Cannot read PPTX files.")
        return None
    except Exception as e:
        logger.error(f"Error reading PPTX {path}: {e}")
        return None


def read_excel(path: Path) -> Optional[str]:
    """Extract text from an Excel file."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_rows = []
            for row in ws.iter_rows(max_row=500, values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    sheet_rows.append(row_text)
            if sheet_rows:
                text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_rows))
        wb.close()

        return "\n\n".join(text_parts).strip() or None
    except ImportError:
        logger.warning("openpyxl not installed. Cannot read Excel files.")
        return None
    except Exception as e:
        logger.error(f"Error reading Excel {path}: {e}")
        return None


def read_csv_file(path: Path) -> Optional[str]:
    """Read CSV file, joining rows into readable text."""
    try:
        raw = path.read_bytes()
        detected = chardet.detect(raw)
        encoding = detected.get("encoding", "utf-8") or "utf-8"
        text = raw.decode(encoding, errors="replace")

        # Parse as CSV for better structure
        reader = csv.reader(io.StringIO(text))
        rows = []
        for i, row in enumerate(reader):
            if i > 500:  # Limit to 500 rows
                rows.append("... (truncated)")
                break
            rows.append(" | ".join(row))

        return "\n".join(rows).strip() or text
    except Exception:
        # Fallback to raw text reading
        return read_text_file(path)


def read_image(path: Path) -> Optional[str]:
    """Extract text from an image using OCR (tesseract)."""
    if not _check_tesseract():
        # Return file metadata as minimal content
        return f"[Image file: {path.name}]"
    try:
        from PIL import Image
        import pytesseract

        img = Image.open(str(path))
        text = pytesseract.image_to_string(img).strip()
        img.close()

        if text:
            return text
        return f"[Image file: {path.name} - no text detected]"
    except ImportError:
        logger.warning("Pillow or pytesseract not installed.")
        return f"[Image file: {path.name}]"
    except Exception as e:
        logger.error(f"Error extracting text from image {path}: {e}")
        return f"[Image file: {path.name}]"


def read_file_content(path: Path) -> Optional[str]:
    """
    Universal file reader. Routes to the appropriate reader based on extension.
    Returns extracted text content, or None if the file can't be read.
    """
    try:
        size = path.stat().st_size
        if size == 0:
            return ""
        if size > MAX_FILE_SIZE_BYTES:
            logger.info(f"Skipping {path} - too large ({size} bytes)")
            return None
    except OSError:
        return None

    suffix = path.suffix.lower()

    # Route to specialized readers for binary formats
    if suffix in PDF_EXTENSIONS:
        return read_pdf(path)
    elif suffix in DOCX_EXTENSIONS:
        return read_docx(path)
    elif suffix in PPTX_EXTENSIONS:
        return read_pptx(path)
    elif suffix in EXCEL_EXTENSIONS:
        return read_excel(path)
    elif suffix in IMAGE_EXTENSIONS:
        return read_image(path)
    elif suffix == ".csv":
        return read_csv_file(path)
    else:
        # Default: treat as text
        return read_text_file(path)
